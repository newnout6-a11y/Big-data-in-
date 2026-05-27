from __future__ import annotations

import hashlib
import json
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

import docx
import pypdf
from qdrant_client.models import (
    Distance,
    FieldCondition,
    Filter,
    FilterSelector,
    MatchValue,
    PayloadSchemaType,
    PointStruct,
    VectorParams,
)

import визуальная_обработка as виз
import извлечение_картинок as извк

_получить_easyocr_reader = извк._получить_easyocr_reader

try:
    import fitz
except ImportError:  # pragma: no cover - optional fallback
    fitz = None


def _ocr_картинки(путь_файла: Path) -> str:
    reader = _получить_easyocr_reader()
    if reader is None:
        return ""
    try:
        import cv2
        import numpy as np
        данные = np.fromfile(str(путь_файла), dtype=np.uint8)
        if данные.size == 0:
            return ""
        изображение = cv2.imdecode(данные, cv2.IMREAD_COLOR)
        if изображение is None:
            return ""
        результаты = reader.readtext(изображение, detail=0, paragraph=True)
    except Exception:
        return ""
    куски = [s.strip() for s in (результаты or []) if s and isinstance(s, str) and s.strip()]
    текст = re.sub(r"\s+", " ", " ".join(куски)).strip()
    return текст[:397].rstrip() + "…" if len(текст) > 400 else текст

# easyocr — опциональная зависимость для распознавания текста на растровых
# схемах. Если не установлена — OCR-функционал просто отключается, всё
# остальное работает.






try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None


# notebooks.py лежит в core/, но user_documents/, extracted_images/ и т. п.
# хранятся в корне репо — поднимаемся на уровень.
BASE_DIR = Path(__file__).resolve().parent.parent
USER_DOCUMENTS_DIR = BASE_DIR / "user_documents"
NOTEBOOKS_FILE = USER_DOCUMENTS_DIR / "notebooks.json"
HIGHLIGHTS_DIR = USER_DOCUMENTS_DIR / "highlights"
EXTRACTED_IMAGES_DIR = BASE_DIR / "extracted_images"

DEFAULT_NOTEBOOKS = ("Физхимия 3 курс", "Диплом", "ML курс Воронцова")
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".pptx"}

VECTOR_SIZE = 768
CHUNK_SIZE = 900
CHUNK_OVERLAP = 120
PSEUDO_PAGE_SIZE = 2500
NOTEBOOK_MIN_SCORE = 0.62


def get_user_id() -> str:
    user_id = (os.getenv("NAVIGATOR_USER_ID") or os.getenv("USER_ID") or "local").strip()
    return _safe_tag(user_id) or "local"


def load_store(user_id: str | None = None) -> dict[str, Any]:
    user_id = user_id or get_user_id()
    USER_DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
    if NOTEBOOKS_FILE.exists():
        try:
            with NOTEBOOKS_FILE.open("r", encoding="utf-8") as f:
                store = json.load(f)
        except Exception:
            store = {}
    else:
        store = {}

    if not isinstance(store, dict):
        store = {}
    store.setdefault("version", 1)
    store.setdefault("users", {})
    user = store["users"].setdefault(user_id, {"notebooks": []})
    user.setdefault("notebooks", [])

    existing_titles = {nb.get("title") for nb in user["notebooks"]}
    changed = False
    for title in DEFAULT_NOTEBOOKS:
        if title not in existing_titles:
            user["notebooks"].append(_new_notebook(title, user_id, stable=True))
            changed = True
    if changed or not NOTEBOOKS_FILE.exists():
        save_store(store)
    return store


def save_store(store: dict[str, Any]) -> None:
    USER_DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
    temp_path = NOTEBOOKS_FILE.with_suffix(".tmp")
    with temp_path.open("w", encoding="utf-8") as f:
        json.dump(store, f, ensure_ascii=False, indent=2)
    temp_path.replace(NOTEBOOKS_FILE)


def list_notebooks(user_id: str | None = None) -> list[dict[str, Any]]:
    user_id = user_id or get_user_id()
    store = load_store(user_id)
    return list(store["users"][user_id]["notebooks"])


def get_notebook(notebook_id: str, user_id: str | None = None) -> dict[str, Any] | None:
    for notebook in list_notebooks(user_id):
        if notebook.get("id") == notebook_id:
            return notebook
    return None


def create_notebook(title: str, user_id: str | None = None) -> dict[str, Any]:
    title = re.sub(r"\s+", " ", title or "").strip()
    if not title:
        raise ValueError("Название тетради пустое.")

    user_id = user_id or get_user_id()
    store = load_store(user_id)
    notebooks = store["users"][user_id]["notebooks"]
    for existing in notebooks:
        if existing.get("title", "").strip().lower() == title.lower():
            return existing

    notebook = _new_notebook(title, user_id, stable=False)
    notebooks.append(notebook)
    save_store(store)
    return notebook


def notebook_label(notebook: dict[str, Any]) -> str:
    count = len(notebook.get("files", []))
    suffix = f" · {count} файл(ов)" if count else " · пусто"
    return f"{notebook.get('title', 'Без названия')}{suffix}"


def ingest_uploaded_files(
    client: Any,
    model: Any,
    notebook_id: str,
    uploads: list[tuple[str, bytes]],
    *,
    user_id: str | None = None,
    visual_mode: bool = False,
    use_groq_vision: bool = False,
    use_ocr: bool = False,
    max_groq_pages_per_file: int = виз.MAX_GROQ_PAGES_DEFAULT,
    on_progress: Callable[[str], None] | None = None,
) -> dict[str, Any]:
    user_id = user_id or get_user_id()
    store = load_store(user_id)
    notebooks = store["users"][user_id]["notebooks"]
    notebook = next((nb for nb in notebooks if nb.get("id") == notebook_id), None)
    if notebook is None:
        raise ValueError("Тетрадь не найдена.")

    ensure_collection(client, notebook)

    summary = {
        "added_files": 0, "skipped_files": 0, "chunks": 0,
        "groq_vision_pages": 0, "ocr_pages": 0, "errors": [],
    }
    known_hashes = {f.get("file_hash") for f in notebook.get("files", [])}

    groq_key = виз._first_groq_key() if use_groq_vision else ""

    for filename, data in uploads:
        safe_name = sanitize_filename(filename)
        ext = Path(safe_name).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            summary["errors"].append(f"{filename}: неподдерживаемый тип файла")
            continue

        file_hash = hashlib.sha256(data).hexdigest()
        if file_hash in known_hashes:
            summary["skipped_files"] += 1
            continue

        target = _file_path(user_id, notebook_id, safe_name, file_hash)
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(data)

        try:
            if visual_mode and ext == ".pdf":
                if on_progress:
                    on_progress(f"{filename}: визуальная обработка...")
                visual_pages = виз.обработать_pdf(
                    target,
                    use_groq_vision=use_groq_vision,
                    groq_api_key=groq_key,
                    max_groq_pages=max_groq_pages_per_file,
                    save_images=True,
                    on_progress=on_progress,
                )
                # Параллельно достаём встроенные диаграммы/картинки. Иначе при
                # включённом visual_mode поиск отдаёт только текст: рендер целой
                # страницы из `визуальная_обработка` нигде в выдаче не выводится.
                images_by_page = _извлечь_встроенные_картинки_по_страницам(
                    target,
                    vision_api_key=groq_key if use_groq_vision else "",
                    use_ocr=use_ocr,
                )
                pages = [
                    {
                        "page": p["page"],
                        "text": p["text"],
                        "images": images_by_page.get(p["page"], []),
                    }
                    for p in visual_pages
                ]
                visual_meta_by_page = {p["page"]: p for p in visual_pages}
                # Страницы, на которых визуальная обработка не дала текста, но
                # есть встроенные картинки — добавляем как «визуальные» чанки,
                # чтобы их можно было поднять через индекс соседних страниц.
                visual_page_numbers = {p["page"] for p in visual_pages}
                for номер_страницы, картинки in images_by_page.items():
                    if номер_страницы in visual_page_numbers or not картинки:
                        continue
                    pages.append({
                        "page": номер_страницы,
                        "text": "",
                        "images": картинки,
                    })
                pages.sort(key=lambda элемент: элемент["page"])
                summary["ocr_pages"] += sum(
                    1 for p in visual_pages if p.get("tier_used") == 1)
                summary["groq_vision_pages"] += sum(
                    1 for p in visual_pages if p.get("tier_used") == 2)
            else:
                pages = extract_pages(target, use_ocr=use_ocr)
                visual_meta_by_page = {}

            chunks = build_chunks(
                pages,
                notebook=notebook,
                user_id=user_id,
                file_hash=file_hash,
                file_path=target,
                original_name=safe_name,
                visual_meta_by_page=visual_meta_by_page,
            )
            if not chunks:
                summary["errors"].append(f"{filename}: не удалось извлечь текст")
                continue
            upsert_chunks(client, model, notebook["collection"], chunks)
        except Exception as error:
            summary["errors"].append(f"{filename}: {error}")
            continue

        notebook.setdefault("files", []).append({
            "name": safe_name,
            "file_hash": file_hash,
            "path": str(target),
            "type": ext.lstrip("."),
            "chunks": len(chunks),
            "uploaded_at": _now_iso(),
            "visual_mode": bool(visual_mode and ext == ".pdf"),
        })
        known_hashes.add(file_hash)
        summary["added_files"] += 1
        summary["chunks"] += len(chunks)

    save_store(store)
    return summary


def переиндексировать_файл_с_ocr(
    client: Any,
    model: Any,
    notebook_id: str,
    file_hash: str,
    *,
    user_id: str | None = None,
    on_progress: Callable[[str], None] | None = None,
) -> dict[str, Any]:
    """Заново индексирует уже загруженный PDF с включённым OCR картинок.

    Поведение:
      1) Удаляет все точки этого файла из Qdrant (по фильтру user/notebook/file_hash).
      2) Перечитывает PDF c `use_ocr=True` — caption'ы картинок-схем
         заполняются распознанным текстом (например «Отстойник
         непрерывного действия с вращающимися скребками…»).
      3) Перестраивает чанки и заливает их в Qdrant.
      4) Обновляет запись файла в `notebooks.json` (счётчик чанков, флаг).

    Сам PDF на диске и каталог `extracted_images/<file_hash>/` не трогаем —
    OCR работает по уже сохранённым там картинкам.
    """
    user_id = user_id or get_user_id()
    store = load_store(user_id)
    notebooks_list = store["users"][user_id]["notebooks"]
    notebook = next((nb for nb in notebooks_list if nb.get("id") == notebook_id), None)
    if notebook is None:
        raise ValueError("Тетрадь не найдена.")
    file_info = next(
        (f for f in notebook.get("files", []) if f.get("file_hash") == file_hash),
        None,
    )
    if file_info is None:
        raise ValueError("Файл не найден в тетради.")

    тип = (file_info.get("type") or "").lower().lstrip(".")
    if тип != "pdf":
        raise ValueError("OCR применим только к PDF-файлам.")

    target = Path(file_info.get("path", ""))
    if not target.exists():
        raise ValueError(f"Исходный PDF не найден на диске: {target}")

    ensure_collection(client, notebook)

    if on_progress:
        on_progress("Удаляю старые точки из Qdrant…")
    try:
        client.delete(
            collection_name=notebook["collection"],
            points_selector=FilterSelector(
                filter=Filter(must=[
                    FieldCondition(key="user_id", match=MatchValue(value=user_id)),
                    FieldCondition(key="notebook_id", match=MatchValue(value=notebook_id)),
                    FieldCondition(key="file_hash", match=MatchValue(value=file_hash)),
                ]),
            ),
        )
    except Exception as ошибка:  # pragma: no cover - сетевые/локальные сбои
        raise ValueError(f"Не удалось удалить старые точки: {ошибка}")

    if on_progress:
        on_progress("Извлекаю текст и применяю OCR к картинкам…")
    pages = _extract_pdf(target, use_ocr=True)
    if not pages:
        raise ValueError("Не удалось извлечь страницы из PDF.")

    chunks = build_chunks(
        pages,
        notebook=notebook,
        user_id=user_id,
        file_hash=file_hash,
        file_path=target,
        original_name=file_info.get("name", target.name),
    )
    if not chunks:
        raise ValueError("После переиндексации не удалось построить чанки.")

    if on_progress:
        on_progress(f"Заливаю {len(chunks)} чанков в Qdrant…")
    upsert_chunks(client, model, notebook["collection"], chunks)

    file_info["chunks"] = len(chunks)
    file_info["reindexed_at"] = _now_iso()
    file_info["use_ocr"] = True
    save_store(store)

    return {"chunks": len(chunks), "pages": len(pages)}


def ensure_collection(client: Any, notebook: dict[str, Any]) -> None:
    collection = notebook["collection"]
    names = {item.name for item in client.get_collections().collections}
    if collection not in names:
        client.create_collection(
            collection_name=collection,
            vectors_config=VectorParams(size=VECTOR_SIZE, distance=Distance.COSINE),
        )
    for field, schema in (
        ("user_id", PayloadSchemaType.KEYWORD),
        ("notebook_id", PayloadSchemaType.KEYWORD),
        ("file_hash", PayloadSchemaType.KEYWORD),
        ("document", PayloadSchemaType.KEYWORD),
        ("source", PayloadSchemaType.KEYWORD),
    ):
        try:
            client.create_payload_index(collection, field, schema)
        except Exception:
            pass


def collection_count(client: Any, notebook: dict[str, Any]) -> int:
    try:
        ensure_collection(client, notebook)
        return int(client.count(notebook["collection"], exact=True).count)
    except Exception:
        return 0


def notebook_documents(notebook: dict[str, Any]) -> list[str]:
    seen = set()
    result = []
    for file_info in notebook.get("files", []):
        name = file_info.get("name")
        if name and name not in seen:
            result.append(name)
            seen.add(name)
    return result


_СТОП_СЛОВА_ПОИСКА: set[str] = {
    "как", "что", "где", "когда", "это", "этот", "эта", "эти", "тот", "та", "те",
    "при", "для", "или", "из", "на", "по", "от", "до", "над", "под", "без", "про",
    "и", "а", "но", "же", "ли", "бы", "не", "ни", "ты", "вы", "он", "она", "они",
    "мы", "я", "его", "её", "их", "сам", "сама", "сами", "свой", "своя",
    "действия", "выглядит", "называется", "такое", "таких", "такая", "такие",
    "the", "a", "an", "is", "are", "of", "in", "on", "to", "for", "with", "and",
    "or", "by", "from", "at", "as", "how", "what", "where", "which", "it",
}


def _стем(слово: str) -> str:
    """Наивный стемминг: для слов длиной ≥6 обрезает до первых 6 символов,
    чтобы «отстойник/отстойника/отстойнике» имели общий корень. Короткие
    слова (3-5 симв) оставляем как есть. Аналог стемминга Портера для
    русского у нас нет без доп. зависимостей — этого приближения хватает
    для keyword-бустинга.
    """
    if len(слово) >= 6:
        return слово[:6]
    return слово


def _токены_запроса(вопрос: str) -> set[str]:
    """Разбивает запрос на нормализованные стем-токены (низкий регистр,
    длина >=3, без стоп-слов). Стемминг даёт совпадение разных словоформ
    («отстойник» в запросе ↔ «отстойнике» в тексте).
    """
    import re as _re
    токены = _re.findall(r"[\wёЁ]+", (вопрос or "").lower(), _re.UNICODE)
    результат: set[str] = set()
    for t in токены:
        if len(t) < 3:
            continue
        if t in _СТОП_СЛОВА_ПОИСКА:
            continue
        результат.add(_стем(t))
    return результат


# Стемы запросов, ориентированных на поиск схем/рисунков/диаграмм.
# Если хотя бы один из этих стемов присутствует в запросе И у чанка есть
# картинки — чанк получает дополнительный буст: предполагается, что страница
# со схемой на изображении релевантнее текстовых страниц, лишь упоминающих
# то же понятие в тексте.
# Префиксы для обнаружения «схемо-ориентированных» запросов.
# Используем startswith вместо точного сравнения стемов: _стем() для слов
# длиной < 6 символов возвращает слово как есть («схема» → «схема», а не «схем»),
# поэтому сравнение с коротким префиксом через startswith надёжнее.
_СХЕМА_ПРЕФИКСЫ = ("схем", "диагр", "рисун", "черт", "граф", "изобр", "слайд")
_IMAGE_БУСТ = 0.12


def _запрос_про_схему(стемы_запроса: set[str]) -> bool:
    """True если запрос ищет схему, рисунок, диаграмму или изображение."""
    return any(т.startswith(p) for т in стемы_запроса for p in _СХЕМА_ПРЕФИКСЫ)


def _keyword_буст(payload: dict[str, Any], стемы_запроса: set[str]) -> float:
    """Возвращает дополнительный скор-буст по числу стем-токенов запроса,
    встречающихся в тексте чанка и caption'ах его картинок.

    Буст = 0.07 за каждый уникальный совпавший стем, с потолком 0.35.
    Плюс +0.12 если запрос содержит «схема»/«диаграмма»/«рисунок» и у
    чанка есть картинки — такой чанк скорее всего содержит искомую схему
    на изображении, а не просто упоминает термин в тексте.
    """
    if not стемы_запроса:
        return 0.0
    import re as _re
    куски: list[str] = []
    текст_чанка = (payload.get("text") or "")
    if текст_чанка:
        куски.append(текст_чанка)
    for картинка in payload.get("images") or []:
        if isinstance(картинка, dict):
            cap = картинка.get("caption") or ""
            if cap:
                куски.append(cap)
    общий = " ".join(куски).lower()
    if not общий:
        return 0.0
    слова_текста = _re.findall(r"[\wёЁ]+", общий, _re.UNICODE)
    стемы_текста = {_стем(w) for w in слова_текста if len(w) >= 3}
    совпало = len(стемы_запроса & стемы_текста)
    if совпало == 0:
        return 0.0
    keyword_score = min(0.35, 0.07 * совпало)
    # Дополнительный буст: запрос ищет схему/рисунок, и у чанка есть
    # картинки — повышаем его шанс обойти текстовые страницы с тем же словом.
    image_score = _IMAGE_БУСТ if (_запрос_про_схему(стемы_запроса) and payload.get("images")) else 0.0
    return keyword_score + image_score


def search_notebook(
    client: Any,
    model: Any,
    notebook: dict[str, Any],
    question: str,
    *,
    limit: int = 5,
    user_id: str | None = None,
    min_score: float = NOTEBOOK_MIN_SCORE,
) -> list[Any]:
    user_id = user_id or get_user_id()
    ensure_collection(client, notebook)
    vector = model.encode("query: " + question, normalize_embeddings=True).tolist()
    query_filter = Filter(must=[
        FieldCondition(key="user_id", match=MatchValue(value=user_id)),
        FieldCondition(key="notebook_id", match=MatchValue(value=notebook["id"])),
    ])
    # Берём глубокий пул кандидатов (limit*6), чтобы до re-rank'а точно
    # дошли страницы с буквальным совпадением ключевых слов — они могут
    # иметь чуть более низкий dense-score, но быть правильным ответом
    # (например, слайд-схема, где ключевой текст только на картинке).
    response = client.query_points(
        collection_name=notebook["collection"],
        query=vector,
        limit=max(limit * 6, limit),
        query_filter=query_filter,
        with_payload=True,
    )
    токены = _токены_запроса(question)
    ранжированные: list[tuple[float, Any]] = []
    for point in response.points:
        payload = point.payload or {}
        text = payload.get("text") or ""
        if not text:
            continue
        dense = float(point.score)
        буст = _keyword_буст(payload, токены)
        итоговый = dense + буст
        # Страницы с реальным keyword-совпадением пропускаем даже при
        # dense-score ниже min_score: без этого слайды-заголовки с
        # критическим текстом только в картинке отсекались бы до rerank'а.
        if dense >= min_score or буст >= 0.15:
            ранжированные.append((итоговый, point))
    ранжированные.sort(key=lambda pair: pair[0], reverse=True)
    return [point for _, point in ранжированные[:limit]]


def notebook_fragments(
    client: Any,
    notebook: dict[str, Any],
    *,
    user_id: str | None = None,
    document: str | None = None,
    limit: int = 80,
) -> list[dict[str, Any]]:
    user_id = user_id or get_user_id()
    ensure_collection(client, notebook)
    must = [
        FieldCondition(key="user_id", match=MatchValue(value=user_id)),
        FieldCondition(key="notebook_id", match=MatchValue(value=notebook["id"])),
    ]
    if document and document != "Все документы":
        must.append(FieldCondition(key="document", match=MatchValue(value=document)))
    query_filter = Filter(must=must)

    points = []
    offset = None
    while len(points) < limit:
        batch_limit = min(64, limit - len(points))
        batch, offset = client.scroll(
            collection_name=notebook["collection"],
            scroll_filter=query_filter,
            limit=batch_limit,
            offset=offset,
            with_payload=True,
            with_vectors=False,
        )
        for point in batch:
            # scroll'ом достаём фрагменты целиком (для UI «все документы тетради»),
            # никакого score из retrieval тут нет — это не поиск, а перечисление.
            points.append(dict(point.payload or {}))
        if offset is None:
            break
    points.sort(key=lambda item: (
        str(item.get("document", "")),
        _int_page(item.get("page")) or 0,
        int(item.get("chunk_index") or 0),
    ))
    return points


def собрать_картинки_по_страницам(
    client: Any,
    notebook: dict[str, Any],
    file_hashes: set[str] | list[str],
    *,
    user_id: str | None = None,
) -> dict[tuple[str, int], list[dict[str, Any]]]:
    """Возвращает индекс `(file_hash, page) -> [images]` по коллекции тетради.

    Скроллит коллекцию тетради с фильтром по user_id/notebook_id, оставляет
    только чанки с непустым `images` и принадлежащие переданным `file_hashes`.
    Дедуп картинок выполняется по полю `path` в пределах одной страницы.
    """
    user_id = user_id or get_user_id()
    file_hashes = set(file_hashes or ())
    if not file_hashes:
        return {}
    ensure_collection(client, notebook)
    query_filter = Filter(must=[
        FieldCondition(key="user_id", match=MatchValue(value=user_id)),
        FieldCondition(key="notebook_id", match=MatchValue(value=notebook["id"])),
    ])
    результат: dict[tuple[str, int], list[dict[str, Any]]] = {}
    видели: dict[tuple[str, int], set[str]] = {}
    offset = None
    while True:
        batch, offset = client.scroll(
            collection_name=notebook["collection"],
            scroll_filter=query_filter,
            limit=128,
            offset=offset,
            with_payload=True,
            with_vectors=False,
        )
        for point in batch:
            payload = point.payload or {}
            fh = payload.get("file_hash")
            if not fh or fh not in file_hashes:
                continue
            картинки = payload.get("images") or []
            if not isinstance(картинки, list) or not картинки:
                continue
            page = _int_page(payload.get("page"))
            if page is None:
                continue
            ключ = (fh, page)
            пути_видели = видели.setdefault(ключ, set())
            список = результат.setdefault(ключ, [])
            for картинка in картинки:
                if not isinstance(картинка, dict):
                    continue
                путь = картинка.get("path") or ""
                if not путь or путь in пути_видели:
                    continue
                пути_видели.add(путь)
                список.append(картинка)
        if offset is None:
            break
    return результат


def extract_pages(path: Path, *, use_ocr: bool = False) -> list[dict[str, Any]]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path, use_ocr=use_ocr)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".txt", ".md"}:
        return _extract_text(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    return []


def build_chunks(
    pages: list[dict[str, Any]],
    *,
    notebook: dict[str, Any],
    user_id: str,
    file_hash: str,
    file_path: Path,
    original_name: str,
    visual_meta_by_page: dict[int, dict[str, Any]] | None = None,
) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    ext = file_path.suffix.lower().lstrip(".")
    visual_meta_by_page = visual_meta_by_page or {}
    for page in pages:
        page_no = page["page"]
        text = clean_text(page["text"])
        page_images = list(page.get("images") or [])
        # Страница без осмысленного текста и без картинок — выкидываем.
        # Раньше выкидывали все страницы с text < 50 символов, и тогда диаграммы
        # с короткой подписью («схема:», «рис. 1») сохранялись на диск, но
        # ни в один чанк не попадали — поиск их не видел.
        if len(text) < 50 and not page_images:
            continue
        meta = visual_meta_by_page.get(page_no, {})
        text_pieces = _split_text(text)
        if not text_pieces and page_images:
            # На странице есть встроенные изображения, но текст слишком короткий
            # для `_split_text`. Эмитим один чанк с тем, что есть, чтобы
            # картинки попали и в свой payload, и в индекс соседних страниц.
            placeholder = text or f"[Страница {page_no}: изображение]"
            text_pieces = [placeholder]
        for chunk_index, chunk_text in enumerate(text_pieces, 1):
            text_hash = hashlib.sha1(chunk_text.encode("utf-8", errors="ignore")).hexdigest()
            chunks.append({
                "text": chunk_text,
                "document": original_name,
                "page": page_no,
                "chunk_index": chunk_index,
                "source": "user_upload",
                "user_id": user_id,
                "notebook_id": notebook["id"],
                "notebook_title": notebook.get("title", ""),
                "file_hash": file_hash,
                "file_path": str(file_path),
                "file_type": ext,
                "title": original_name,
                "text_hash": text_hash,
                "uploaded_at": _now_iso(),
                "tier_used": meta.get("tier_used", 0),
                "has_ocr": bool(meta.get("has_ocr", False)),
                "has_visual_caption": bool(meta.get("has_visual_caption", False)),
                "page_hash": meta.get("page_hash", ""),
                "page_image_path": meta.get("image_path", ""),
                "images": page_images,
            })
    return chunks


def _текст_для_эмбеддинга(chunk: dict[str, Any]) -> str:
    """Объединяет текст чанка с caption'ами его картинок. Так OCR-текст со
    схем (например «Отстойник непрерывного действия с вращающимися
    скребками») попадает в эмбеддинг и поиск находит страницу даже если
    этих слов нет в самом текстовом слое PDF. Сам `chunk['text']` остаётся
    для отображения.

    Caption'ы встраиваются как естественное продолжение текста слайда, без
    служебного префикса «Изображения на странице:». Префикс даёт модели
    сигнал «это метаданные, можно недовешивать», из-за чего короткие
    текстовые слайды с критически важным текстом ТОЛЬКО в картинке
    проигрывали ранжирование длинным текстам с совпадающими словами
    (например, слайду про экстракцию, где «отстойник» просто упоминается).
    """
    база = (chunk.get("text") or "").strip()
    captions: list[str] = []
    видели: set[str] = set()
    for картинка in chunk.get("images") or []:
        cap = (картинка.get("caption") or "").strip()
        if not cap:
            continue
        ключ = cap.lower()
        if ключ in видели:
            continue
        видели.add(ключ)
        captions.append(cap)
    if not captions:
        return база
    хвост = "\n".join(captions)
    if база:
        return f"{база}\n\n{хвост}"
    return хвост


def upsert_chunks(client: Any, model: Any, collection: str, chunks: list[dict[str, Any]]) -> None:
    texts = [_текст_для_эмбеддинга(chunk) for chunk in chunks]
    vectors = model.encode(
        ["passage: " + text for text in texts],
        normalize_embeddings=True,
        show_progress_bar=False,
    )
    points = []
    for chunk, vector in zip(chunks, vectors):
        point_id = str(uuid.uuid5(
            uuid.NAMESPACE_URL,
            f"{chunk['notebook_id']}:{chunk['file_hash']}:{chunk['page']}:{chunk['chunk_index']}:{chunk['text_hash']}",
        ))
        points.append(PointStruct(id=point_id, vector=vector.tolist(), payload=chunk))
    client.upsert(collection_name=collection, points=points)


def citation_url(fragment: Any) -> str | None:
    payload = _payload(fragment)
    target = source_file_path(payload)
    if target is None:
        return None

    page = _int_page(payload.get("page"))
    if target.suffix.lower() == ".pdf" and page:
        highlighted = _highlighted_pdf(target, page, payload.get("text", ""))
        if highlighted and highlighted.exists():
            target = highlighted
        return target.resolve().as_uri() + f"#page={page}"
    return target.resolve().as_uri()


def source_file_path(fragment: Any) -> Path | None:
    payload = _payload(fragment)
    path = payload.get("file_path")
    document = payload.get("document") or ""
    candidates = []
    if path:
        candidates.append(Path(path))
    if document:
        candidates.extend([
            BASE_DIR / "all_pdfs" / document,
            BASE_DIR / "user_documents" / document,
        ])
    for candidate in candidates:
        try:
            if candidate.exists() and candidate.is_file():
                return candidate
        except OSError:
            continue
    return None


def downloadable_name(fragment: Any) -> str:
    payload = _payload(fragment)
    path = source_file_path(payload)
    if path is not None:
        return path.name
    return sanitize_filename(payload.get("document") or "document")
































def _extract_pdf(
    path: Path,
    *,
    vision_api_key: str = "",
    use_ocr: bool = False,
) -> list[dict[str, Any]]:
    """Возвращает список страниц `{page, text, images}`.

    `images` — это список встроенных в PDF картинок, сохранённых на диск под
    `extracted_images/<file_hash>/`. Извлекаются всегда, даже без флага
    `visual_mode` — это отдельная (бесплатная) фича: «показать рисунки рядом
    с найденным фрагментом». Если задан `vision_api_key`, для картинок без
    подписи в тексте вызывается Groq Vision.
    """
    pages: list[dict[str, Any]] = []
    file_hash = _file_hash(path)
    images_dir = EXTRACTED_IMAGES_DIR / file_hash
    if fitz is not None:
        try:
            doc = fitz.open(path)
            try:
                фоновые = извк.найти_фоновые_xref(doc)
                for index, page in enumerate(doc, start=1):
                    text = page.get_text("text") or ""
                    images = извк.извлечь_картинки_страницы(
                        doc, page, index, images_dir,
                        текст_страницы=text,
                        vision_api_key=vision_api_key,
                        фоновые_xref=фоновые,
                        use_ocr=use_ocr,
                        base_dir=BASE_DIR,
                    )
                    if len(text.strip()) > 40 or images:
                        pages.append({"page": index, "text": text, "images": images})
            finally:
                doc.close()
        except Exception:
            pages = []
    if pages:
        return pages

    try:
        reader = pypdf.PdfReader(str(path))
        for index, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            if len(text.strip()) > 40:
                pages.append({"page": index, "text": text, "images": []})
    except Exception:
        return []
    return pages


def _извлечь_встроенные_картинки_по_страницам(
    path: Path,
    *,
    vision_api_key: str = "",
    use_ocr: bool = False,
) -> dict[int, list[dict[str, Any]]]:
    """Сохраняет встроенные в PDF картинки и возвращает индекс `page -> [images]`.

    Используется в `ingest_uploaded_files`, когда `visual_mode=True`: визуальная
    обработка возвращает только текст, поэтому без отдельного прохода по PDF
    встроенные диаграммы не попадают в payload чанков, и поиск отдаёт только
    текст. Формат элементов совпадает с `_извлечь_картинки_страницы`.
    """
    if fitz is None:
        return {}
    file_hash = _file_hash(path)
    images_dir = EXTRACTED_IMAGES_DIR / file_hash
    индекс: dict[int, list[dict[str, Any]]] = {}
    try:
        doc = fitz.open(path)
    except Exception:
        return {}
    try:
        фоновые = извк.найти_фоновые_xref(doc)
        for index, page in enumerate(doc, start=1):
            текст = ""
            try:
                текст = page.get_text("text") or ""
            except Exception:
                текст = ""
            картинки = извк.извлечь_картинки_страницы(
                doc, page, index, images_dir,
                текст_страницы=текст,
                vision_api_key=vision_api_key,
                фоновые_xref=фоновые,
                use_ocr=use_ocr,
                base_dir=BASE_DIR,
            )
            if картинки:
                индекс[index] = картинки
    finally:
        doc.close()
    return индекс


def _file_hash(path: Path) -> str:
    """SHA-256 от содержимого файла, поточно. Используется как подкаталог для
    извлечённых картинок и для дедупа в payload (поле file_hash).

    Раньше тут был SHA-1 — несогласованный с другими местами проекта,
    которые считают SHA-256. Старые папки `extracted_images/<sha1>/` остаются
    на диске и продолжают работать через сохранённые в payload пути; новые
    загрузки идут уже под SHA-256.
    """
    h = hashlib.sha256()
    try:
        with path.open("rb") as f:
            for блок in iter(lambda: f.read(65536), b""):
                h.update(блок)
    except OSError:
        return "unknown"
    return h.hexdigest()


def _extract_docx(path: Path) -> list[dict[str, Any]]:
    document = docx.Document(str(path))
    pieces = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                pieces.append(row_text)
    return _pseudo_pages("\n".join(pieces))


def _extract_text(path: Path) -> list[dict[str, Any]]:
    data = path.read_bytes()
    for encoding in ("utf-8", "cp1251", "latin-1"):
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = data.decode("utf-8", errors="ignore")
    return _pseudo_pages(text)


def _extract_pptx(path: Path) -> list[dict[str, Any]]:
    if Presentation is None:
        raise RuntimeError("Для PPTX нужен пакет python-pptx.")
    deck = Presentation(str(path))
    pages = []
    for index, slide in enumerate(deck.slides, start=1):
        pieces = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                pieces.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        pieces.append(row_text)
        text = "\n".join(piece for piece in pieces if piece)
        if len(text.strip()) > 20:
            pages.append({"page": index, "text": text})
    return pages


def _pseudo_pages(text: str) -> list[dict[str, Any]]:
    text = clean_text(text)
    pages = []
    start = 0
    page = 1
    while start < len(text):
        chunk = text[start:start + PSEUDO_PAGE_SIZE].strip()
        if len(chunk) > 40:
            pages.append({"page": page, "text": chunk})
        start += PSEUDO_PAGE_SIZE
        page += 1
    return pages


def _split_text(text: str) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + CHUNK_SIZE)
        piece = text[start:end].strip()
        if len(piece) > 80:
            chunks.append(piece)
        if end >= len(text):
            break
        start = max(0, end - CHUNK_OVERLAP)
    return chunks


def clean_text(text: str) -> str:
    text = text.replace("\u00ad", "")
    text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def sanitize_filename(name: str) -> str:
    name = Path(name or "document").name
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name)
    name = re.sub(r"\s+", " ", name).strip(" .")
    if not name:
        name = "document"
    stem = Path(name).stem[:90] or "document"
    suffix = Path(name).suffix.lower()
    return stem + suffix


def _file_path(user_id: str, notebook_id: str, filename: str, file_hash: str) -> Path:
    return USER_DOCUMENTS_DIR / user_id / notebook_id / "files" / f"{file_hash[:12]}_{filename}"


def _new_notebook(title: str, user_id: str, *, stable: bool) -> dict[str, Any]:
    seed = f"{user_id}:{title}" if stable else f"{user_id}:{title}:{_now_iso()}"
    short_hash = hashlib.sha1(seed.encode("utf-8")).hexdigest()[:12]
    notebook_id = f"nb_{short_hash}"
    return {
        "id": notebook_id,
        "title": title,
        "collection": f"user_nb_{short_hash}",
        "created_at": _now_iso(),
        "files": [],
    }


def _safe_tag(value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9_.-]+", "_", value.strip())
    return value[:80].strip("_.-")


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _payload(fragment: Any) -> dict[str, Any]:
    if isinstance(fragment, dict):
        return fragment
    return getattr(fragment, "payload", {}) or {}


def _int_page(value: Any) -> int | None:
    try:
        page = int(value)
    except (TypeError, ValueError):
        return None
    return page if page > 0 else None


def _highlighted_pdf(path: Path, page_number: int, quote: str) -> Path | None:
    if fitz is None:
        return path
    HIGHLIGHTS_DIR.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha1(
        f"{path.resolve()}:{page_number}:{quote[:1200]}".encode("utf-8", errors="ignore")
    ).hexdigest()[:16]
    safe_stem = re.sub(r"[^a-zA-Z0-9_.-]+", "_", path.stem)[:70] or "document"
    output = HIGHLIGHTS_DIR / f"{safe_stem}__p{page_number}__{digest}.pdf"
    if output.exists():
        return output

    try:
        doc = fitz.open(path)
        try:
            if page_number > len(doc):
                return path
            page = doc[page_number - 1]
            rect = _best_text_rect(page, quote)
            if rect is not None:
                annot = page.add_highlight_annot(rect)
                try:
                    annot.set_colors(stroke=(1, 0.82, 0.15))
                    annot.update(opacity=0.35)
                except Exception:
                    annot.update()
                doc.save(output, garbage=4, deflate=True)
                return output
        finally:
            doc.close()
    except Exception:
        return path
    return path


def _best_text_rect(page: Any, quote: str) -> Any | None:
    clean_quote = re.sub(r"\s+", " ", quote or "").strip()
    if not clean_quote:
        return None

    for snippet in _search_snippets(clean_quote):
        try:
            rects = page.search_for(snippet)
        except Exception:
            rects = []
        if rects:
            return rects[0]

    quote_words = set(re.findall(r"[A-Za-zА-Яа-яЁё0-9]{4,}", clean_quote.lower()))
    if not quote_words:
        return None
    best_rect = None
    best_score = 0.0
    try:
        blocks = page.get_text("blocks") or []
    except Exception:
        blocks = []
    for block in blocks:
        if len(block) < 5:
            continue
        block_text = str(block[4])
        block_words = set(re.findall(r"[A-Za-zА-Яа-яЁё0-9]{4,}", block_text.lower()))
        if not block_words:
            continue
        overlap = len(quote_words & block_words)
        score = overlap / max(1, min(len(quote_words), len(block_words)))
        if score > best_score:
            best_score = score
            best_rect = fitz.Rect(block[:4]) if fitz is not None else None
    return best_rect if best_score >= 0.12 else None


def _search_snippets(text: str) -> list[str]:
    snippets = []
    for candidate in re.split(r"(?<=[.!?])\s+", text):
        candidate = candidate.strip()
        if 45 <= len(candidate) <= 180:
            snippets.append(candidate)
            break
    for size in (160, 120, 80):
        if len(text) >= size:
            snippets.append(text[:size].strip())
    return list(dict.fromkeys(snippets))
